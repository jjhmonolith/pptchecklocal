#!/usr/bin/env python3
"""
PowerPoint 파일 분석을 위한 Python 스크립트
python-pptx 라이브러리를 사용하여 PPTX 파일의 텍스트를 추출하고 구조화된 데이터로 반환
"""

import sys
import json
import requests
from io import BytesIO
import tempfile
import os

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def download_file(url, temp_path):
    """URL에서 파일 다운로드"""
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        with open(temp_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        return True
    except Exception as e:
        print(f"파일 다운로드 실패: {str(e)}", file=sys.stderr)
        return False

def extract_text_from_shape(shape):
    """도형에서 텍스트 추출"""
    text_runs = []
    
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return text_runs
        
    for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        for run_idx, run in enumerate(paragraph.runs):
            if run.text.strip():  # 빈 문자열이 아닌 경우만
                text_runs.append({
                    "text": run.text,
                    "paragraph_idx": paragraph_idx,
                    "run_idx": run_idx,
                    "font_name": run.font.name if run.font.name else "기본",
                    "font_size": run.font.size.pt if run.font.size else 12,
                    "is_bold": run.font.bold if run.font.bold is not None else False,
                    "is_italic": run.font.italic if run.font.italic is not None else False
                })
    
    return text_runs

def analyze_pptx_file(file_path):
    """PPTX 파일 분석"""
    try:
        # PowerPoint 프레젠테이션 로드
        prs = Presentation(file_path)
        
        result = {
            "slides": [],
            "stats": {
                "slides": len(prs.slides),
                "shapes": 0,
                "runs": 0,
                "tokensEstimated": 0
            }
        }
        
        # 각 슬라이드 처리
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                "slideIndex": slide_idx + 1,
                "shapes": []
            }
            
            # 각 도형 처리
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or \
                   shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or \
                   (hasattr(shape, 'text_frame') and shape.text_frame):
                    
                    text_runs = extract_text_from_shape(shape)
                    
                    if text_runs:  # 텍스트가 있는 경우만
                        shape_data = {
                            "shapeId": f"shape-{slide_idx + 1}-{shape_idx + 1}",
                            "shapeType": str(shape.shape_type),
                            "textRuns": text_runs
                        }
                        
                        slide_data["shapes"].append(shape_data)
                        result["stats"]["shapes"] += 1
                        result["stats"]["runs"] += len(text_runs)
                        
                        # 토큰 수 추정 (한글 기준 대략적 계산)
                        total_text = "".join([run["text"] for run in text_runs])
                        result["stats"]["tokensEstimated"] += len(total_text.replace(" ", "")) * 1.5
            
            if slide_data["shapes"]:  # 텍스트가 있는 슬라이드만 추가
                result["slides"].append(slide_data)
        
        return result
        
    except Exception as e:
        raise Exception(f"PPTX 파일 분석 오류: {str(e)}")

def main():
    """메인 함수"""
    if len(sys.argv) != 2:
        print(json.dumps({
            "error": "사용법: python pptx_analyzer.py <file_url>"
        }))
        sys.exit(1)
    
    file_url = sys.argv[1]
    
    # python-pptx 라이브러리 확인
    if not PPTX_AVAILABLE:
        print(json.dumps({
            "error": "python-pptx 라이브러리가 설치되지 않았습니다. 'pip install python-pptx requests' 명령을 실행하세요."
        }))
        sys.exit(1)
    
    # 임시 파일 생성
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        temp_path = temp_file.name
    
    try:
        # 파일 다운로드
        if not download_file(file_url, temp_path):
            print(json.dumps({
                "error": "파일 다운로드에 실패했습니다."
            }))
            sys.exit(1)
        
        # PPTX 파일 분석
        result = analyze_pptx_file(temp_path)
        
        # JSON 형태로 결과 출력
        print(json.dumps(result, ensure_ascii=False, indent=2))
        
    except Exception as e:
        print(json.dumps({
            "error": str(e)
        }), file=sys.stderr)
        sys.exit(1)
        
    finally:
        # 임시 파일 정리
        try:
            os.unlink(temp_path)
        except:
            pass

if __name__ == "__main__":
    main()